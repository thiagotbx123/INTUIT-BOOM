"""
RACCOONS PPTs v2 - Visual com Cards e Hyperlinks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# Cores
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x41, 0x69, 0xE1)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x28, 0xA7, 0x45)
LIGHT_GREEN = RGBColor(0xD4, 0xED, 0xDA)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
LIGHT_ORANGE = RGBColor(0xFF, 0xEC, 0xD9)
PURPLE = RGBColor(0x6F, 0x42, 0xC1)
LIGHT_PURPLE = RGBColor(0xE9, 0xD5, 0xFF)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)


def add_hyperlink(run, url):
    """Adiciona hyperlink a um run de texto"""
    rId = run.part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    rPr = run._r.get_or_add_rPr()
    hlinkClick = parse_xml(
        f'<a:hlinkClick xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:id="{rId}"/>'
    )
    rPr.append(hlinkClick)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient effect (solid dark)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(9), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER


def add_card_grid_slide(prs, title, cards):
    """Slide com grid de cards coloridos e hyperlinks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Cards em grid 2x3
    colors = [
        (LIGHT_BLUE, ACCENT_BLUE),
        (LIGHT_GREEN, GREEN),
        (LIGHT_ORANGE, ORANGE),
        (LIGHT_PURPLE, PURPLE),
        (LIGHT_GRAY, GRAY),
        (LIGHT_BLUE, ACCENT_BLUE),
    ]

    for i, (card_title, card_desc, card_url) in enumerate(cards):
        col = i % 2
        row = i // 2

        x = Inches(0.4 + col * 4.8)
        y = Inches(1.4 + row * 2.0)
        w = Inches(4.5)
        h = Inches(1.8)

        bg_color, accent_color = colors[i % len(colors)]

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = accent_color
        card.line.width = Pt(2)

        # Accent bar no topo do card
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()

        # Numero do card
        num_box = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.2), Inches(0.5), Inches(0.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = accent_color

        # Titulo do card
        title_box = slide.shapes.add_textbox(
            x + Inches(0.6), y + Inches(0.25), Inches(3.7), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Descricao
        desc_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.75), Inches(4.1), Inches(0.5)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

        # Link clicavel
        if card_url:
            link_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(1.35), Inches(4.1), Inches(0.4)
            )
            tf = link_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Abrir >>"
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = accent_color
            run.font.underline = True
            add_hyperlink(run, card_url)


def add_diagram_slide(prs, title, items):
    """Slide com diagrama de fluxo horizontal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Diagrama horizontal
    n = len(items)
    box_width = Inches(2.8)
    spacing = (prs.slide_width - box_width * min(n, 3)) / (min(n, 3) + 1)

    for i, (item_title, item_desc) in enumerate(items):
        row = i // 3
        col = i % 3

        x = spacing + col * (box_width + Inches(0.2))
        y = Inches(1.6 + row * 2.8)

        # Box principal
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, Inches(2.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(2)

        # Numero circular
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.1), y - Inches(0.25), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(
            x + Inches(1.1), y - Inches(0.2), Inches(0.5), Inches(0.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Titulo
        t_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.4), Inches(2.6), Inches(0.6)
        )
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

        # Descricao
        d_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(1.0), Inches(2.6), Inches(1.2)
        )
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item_desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER


def add_action_cards_slide(prs, title, actions):
    """Slide com cards de acoes lado a lado"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header verde
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Duas colunas
    col_width = Inches(4.6)

    for i, (action, owner) in enumerate(actions):
        col = i % 2
        row = i // 2

        x = Inches(0.3) + col * Inches(4.9)
        y = Inches(1.4) + row * Inches(1.0)

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, Inches(0.85)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GREEN if col == 0 else LIGHT_BLUE
        card.line.color.rgb = GREEN if col == 0 else ACCENT_BLUE
        card.line.width = Pt(1.5)

        # Checkbox visual
        check = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.1),
            y + Inches(0.25),
            Inches(0.3),
            Inches(0.3),
        )
        check.fill.solid()
        check.fill.fore_color.rgb = WHITE
        check.line.color.rgb = GREEN if col == 0 else ACCENT_BLUE

        # Texto acao
        t_box = slide.shapes.add_textbox(
            x + Inches(0.5), y + Inches(0.15), Inches(3.0), Inches(0.55)
        )
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Owner badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(3.6),
            y + Inches(0.25),
            Inches(0.9),
            Inches(0.35),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN if col == 0 else ACCENT_BLUE
        badge.line.fill.background()

        o_box = slide.shapes.add_textbox(
            x + Inches(3.6), y + Inches(0.28), Inches(0.9), Inches(0.35)
        )
        tf = o_box.text_frame
        p = tf.paragraphs[0]
        p.text = owner
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


# ============================================
# PPT 1: RECURSOS E LINKS
# ============================================
prs1 = Presentation()
prs1.slide_width = Inches(10)
prs1.slide_height = Inches(7.5)

add_title_slide(prs1, "RACCOONS", "Recursos e Links Centrais")

# Cards com hyperlinks
recursos = [
    (
        "Control Center",
        "Hub central de gestao do time",
        "https://drive.google.com/drive/folders/1j29MjAzZzKGySd7IBRauqkHKlUU-CaBT",
    ),
    (
        "Work Log (Linear)",
        "Registro de atividades - RAC-16",
        "https://linear.app/testbox/issue/RAC-16/worklog-251217-tsa-thiago-rodrigues",
    ),
    (
        "Client Overview",
        "Visao consolidada de clientes (Coda)",
        "https://coda.io/d/Solutions-Central_djfymaxsTtA/Client-Overview_suQxM0kT#_lue-8nKH",
    ),
    (
        "Planilha Ferias",
        "Controle de ausencias do time",
        "https://docs.google.com/spreadsheets/d/1hWB5mz0upnRJtSxM9Qp4iYn-v2V-Z584r35wSD60XD0/edit?gid=304700225#gid=304700225",
    ),
    (
        "Job Description",
        "Descricoes de cargo e competencias",
        "https://docs.google.com/spreadsheets/d/1SocOzWQAUWpv3BH5OkEdKRp-SSs6k-Cp/edit?gid=276597511#gid=276597511",
    ),
    ("Daily Report", "Consolidacao e padrao de preenchimento", ""),
]
add_card_grid_slide(prs1, "Recursos Centrais", recursos)

# Projetos de melhoria como diagrama
projetos = [
    ("CORTEX", "Sistema de coleta e consolidacao de conhecimento entre fontes"),
    ("SpineHUB", "Backbone de persistencia e memoria entre sessoes"),
    ("Health Checker", "Validacao automatizada de features QBO"),
]
add_diagram_slide(prs1, "Projetos de Melhoria", projetos)

prs1.save("C:/Users/adm_r/intuit-boom/RACCOONS_Recursos_Links.pptx")
print("[OK] PPT 1 atualizado: RACCOONS_Recursos_Links.pptx")

# ============================================
# PPT 2: ALINHAMENTOS E PROXIMOS PASSOS
# ============================================
prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)

add_title_slide(prs2, "SYNC RACCOONS", "Alinhamentos & Proximos Passos | 02 Jan 2025")

# Alinhamentos como diagrama de processo
alinhamentos = [
    (
        "Review Data Eng",
        "Demandas com geracao de dados, IA ou mudanca em dataset passam por review",
    ),
    (
        "Visibilidade Thais",
        "Criterio de demanda na daily dos TSAs para filtrar Data Gen",
    ),
    ("Criterios Formais", "Thais vai formalizar o que entra como 'geracao de dados'"),
    ("Prazos Integrados", "Envolver Thais na definicao de prazos de SOWs e entregas"),
    ("Todas as Areas", "TSAs consideram datas de TODAS areas que contribuem"),
    ("Licoes Aprendidas", "Feedbacks periodicos - Top 10 erros para melhoria continua"),
]
add_diagram_slide(prs2, "Alinhamentos com Thais", alinhamentos)

# Novas iniciativas
iniciativas = [
    ("Prompting TSAs", "Refinamento do treinamento de prompting para o time"),
    (
        "Consultoria Prompts",
        "Nova categoria: consultoria em projetos internos (Cortex/SpineHUB)",
    ),
    ("SYNC Periodico", "Periodicidade definida + script padrao de agenda"),
]
add_diagram_slide(prs2, "Novas Iniciativas", iniciativas)

# Acoes com owners
acoes = [
    ("Criar criterio Data Gen na daily", "Thiago"),
    ("Formalizar criterios 'geracao dados'", "Thais"),
    ("Definir quando review obrigatorio", "Thais"),
    ("Incluir projetos no daily report", "Thiago"),
    ("Alinhar consultoria prompts", "Thais"),
    ("Cascatear alinhamentos TSAs", "Thiago"),
]
add_action_cards_slide(prs2, "Proximos Passos", acoes)

prs2.save("C:/Users/adm_r/intuit-boom/RACCOONS_Alinhamentos_Sync.pptx")
print("[OK] PPT 2 atualizado: RACCOONS_Alinhamentos_Sync.pptx")

print("")
print("=" * 50)
print("2 PPTs atualizados com hyperlinks e visual!")
print("=" * 50)
