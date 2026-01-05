"""
Cria 2 PPTs para reunião do time RACCOONS
- PPT 1: Recursos e Links Centrais
- PPT 2: Alinhamentos Thais + Próximos Passos
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cores TestBox
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x41, 0x69, 0xE1)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)


def add_title_slide(prs, title, subtitle=""):
    """Slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = ACCENT_BLUE
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, items, use_icons=True):
    """Slide com lista de itens"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Items
    y_pos = 1.5
    icons = ["[1]", "[2]", "[3]", "[4]", "[5]", "[6]", "[7]", "[8]"]

    for i, (item_title, item_desc) in enumerate(items):
        # Box para cada item
        item_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8)
        )
        tf = item_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        icon = icons[i % len(icons)] if use_icons else "•"
        p.text = f"{icon}  {item_title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        if item_desc:
            p2 = tf.add_paragraph()
            p2.text = f"     {item_desc}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        y_pos += 0.75

    return slide


def add_bullet_slide(prs, title, bullets):
    """Slide com bullets detalhados"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f">  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(12)

    return slide


def add_action_slide(prs, title, actions):
    """Slide de ações com owners"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_BLUE
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Ações
    y_pos = 1.5
    for action, owner in actions:
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(y_pos),
            Inches(9),
            Inches(0.7),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()

        # Texto
        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.15), Inches(6.5), Inches(0.5)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Owner badge
        owner_box = slide.shapes.add_textbox(
            Inches(7.5), Inches(y_pos + 0.15), Inches(1.8), Inches(0.5)
        )
        tf = owner_box.text_frame
        p = tf.paragraphs[0]
        p.text = owner
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.RIGHT

        y_pos += 0.85

    return slide


# ============================================
# PPT 1: RECURSOS E LINKS DO TIME
# ============================================
prs1 = Presentation()
prs1.slide_width = Inches(10)
prs1.slide_height = Inches(7.5)

add_title_slide(prs1, "RACCOONS", "Recursos e Links Centrais")

recursos = [
    ("Control Center", "Hub central de gestão do time"),
    ("Work Log (Linear)", "Registro de atividades diárias - RAC-16"),
    ("Client Overview (Coda)", "Visão consolidada de clientes"),
    ("Planilha Férias", "Controle de ausências do time"),
    ("Job Description", "Descrições de cargo e competências"),
    ("Daily Report", "Consolidação e padrão de preenchimento"),
]
add_content_slide(prs1, "Recursos Centrais", recursos)

projetos = [
    ("Cortex", "Sistema de coleta e consolidação de conhecimento"),
    ("SpineHUB", "Backbone de persistência e memória"),
    ("Intuit Health Checker", "Validação automatizada de features QBO"),
]
add_content_slide(prs1, "Projetos de Melhoria", projetos)

prs1.save("C:/Users/adm_r/intuit-boom/RACCOONS_Recursos_Links.pptx")
print("[OK] PPT 1 criado: RACCOONS_Recursos_Links.pptx")

# ============================================
# PPT 2: ALINHAMENTOS E PRÓXIMOS PASSOS
# ============================================
prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)

add_title_slide(prs2, "SYNC RACCOONS", "Alinhamentos & Próximos Passos\n02 Jan 2025")

# Alinhamentos Thais
alinhamentos = [
    "Demandas com geração de dados, IA ou mudança em dataset → Review obrigatório Data Eng",
    "Thais terá visibilidade via critério de demanda na daily dos TSAs",
    "Thais vai formalizar critérios do que entra como 'geração de dados'",
    "Envolver Thais na definição de prazos de SOWs e entregas",
    "TSAs devem considerar datas de TODAS áreas que contribuem para entrega",
    "Thais trará feedbacks periódicos de lições aprendidas (ex: top 10 erros)",
]
add_bullet_slide(prs2, "Alinhamentos com Thais", alinhamentos)

# Novas iniciativas
iniciativas = [
    "Refinamento do treinamento de prompting para TSAs",
    "Nova categoria: Consultoria em refinamento de prompts de projetos internos",
    "Cascatear alinhamentos relevantes para time de TSAs",
    "Periodicidade definida para SYNC + script padrão de agenda",
]
add_bullet_slide(prs2, "Novas Iniciativas", iniciativas)

# Ações
acoes = [
    ("Criar critério de demanda Data Gen na daily", "Thiago"),
    ("Formalizar critérios de 'geração de dados'", "Thais"),
    ("Definir quando review é obrigatório", "Thais"),
    ("Incluir projetos de melhoria no daily report", "Thiago"),
    ("Alinhar consultoria prompts (Cortex/SpineHUB)", "Thais"),
    ("Cascatear alinhamentos para TSAs", "Thiago"),
]
add_action_slide(prs2, "Próximos Passos", acoes)

prs2.save("C:/Users/adm_r/intuit-boom/RACCOONS_Alinhamentos_Sync.pptx")
print("[OK] PPT 2 criado: RACCOONS_Alinhamentos_Sync.pptx")

print("")
print("=" * 50)
print("2 PPTs criados com sucesso!")
print("=" * 50)
